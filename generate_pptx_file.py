from detect_price_change import detect_price_change
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

price_changes = detect_price_change()

prs = Presentation()
slide_layout = prs.slide_layouts[6]

for change in price_changes:
    product_name = change['name']
    product_price = f"₦{change['new_rate']}"

    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    logo_path = "logo.png"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(1.2))

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(3), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "your padi padi store"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Inches(4.5), Inches(0.6), Inches(4), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PRICE LIST"
    p.font.name = "Segoe Print"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = product_name
    p.font.name = "Segoe Print"
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = product_price
    p.font.name = "Segoe Print"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

output_path = 'price_list_slide_all.pptx'
prs.save(output_path)
os.startfile(output_path)