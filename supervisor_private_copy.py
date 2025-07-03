"""
Generates a detailed PowerPoint report for supervisors, 
focusing on product price changes and critical anomalies.
This report is intended only for internal review and decision-making.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os


def create_pptx():
    with open("price_changes.json", 'r', encoding='utf-8') as f:
        data = json.load(f)
    threshold = 30
    prs = Presentation()

    
    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
    summary_slide.shapes.title.text = "Price Change Summary Dashboard"
    for paragraph in summary_slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Segoe Print"

    total_changes = len(data)
    anomalies = [item for item in data if abs(item['percent_change']) > threshold]

    summary = f"""
    Total Price Changes: {total_changes}
    Anomalies Detected (> {threshold}% change): {len(anomalies)}
    """
    textbox = summary_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    tf = textbox.text_frame
    tf.text = summary.strip()
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Segoe Print"

    
    detail_slide = prs.slides.add_slide(prs.slide_layouts[5])
    detail_slide.shapes.title.text = "Top 5 Biggest Price Changes"
    for paragraph in detail_slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Segoe Print"

    sorted_changes = sorted(data, key=lambda x: abs(x['percent_change']), reverse=True)[:5]
    text = ""
    for item in sorted_changes:
        text += (
            f"{item['name']} (SKU: {item['sku']}): {item['old_rate']} → {item['new_rate']} "
            f"({item['percent_change']}%)\n"
        )

    textbox = detail_slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(4))
    tf = textbox.text_frame
    tf.text = text.strip()
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Segoe Print"

    
    anomaly_slide = prs.slides.add_slide(prs.slide_layouts[5])
    anomaly_slide.shapes.title.text = "Price Anomalies (Critical Changes)"
    for paragraph in anomaly_slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Segoe Print"

    if anomalies:
        text = ""
        for item in anomalies:
            text += (
                f"{item['name']} (SKU: {item['sku']}): {item['old_rate']} → {item['new_rate']} "
                f"({item['percent_change']}%)\n"
            )
    else:
        text = "No anomalies detected."

    textbox = anomaly_slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(4))
    tf = textbox.text_frame
    tf.text = text.strip()
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Segoe Print"

    prs.save("supervisor_price_dashboard.pptx")
    os.startfile("supervisor_price_dashboard.pptx")


if __name__ == "__main__":
    create_pptx()